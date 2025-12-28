
# -*- coding: utf-8 -*-
"""
api_indexacao.py
- Indexa conteúdo de arquivos (files) em FTS5 (docs) com controle incremental via 'map'.
- Suporta: PDF, DOCX, PPTX, XLSX, TXT, CSV.
"""

import csv
from typing import Optional, List
from fastapi import APIRouter, Depends, HTTPException, Query
from pydantic import BaseModel
from sqlalchemy import text
from sqlalchemy.orm import Session

from app.models.models import RootFolder, File
from app.db.database import get_db
from app.core.deps import require_root_access

# Extratores (locais)
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

router = APIRouter(prefix="/index", tags=["Indexação"])

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".csv"}
BATCH_SIZE = 200

# --------- utilidades de extração ---------
def safe_read_text(path, encodings=("utf-8", "latin-1", "cp1252")):
    for enc in encodings:
        try:
            with open(path, "r", encoding=enc, errors="ignore") as f:
                return f.read()
        except Exception:
            continue
    return ""

def extract_pdf(path):
    try:
        reader = PdfReader(path)
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text() or "")
        return "\n".join(texts)
    except Exception:
        return ""

def extract_docx(path):
    try:
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception:
        return ""

def extract_pptx(path):
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""

def extract_xlsx(path, max_cells=10000):
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        texts = []
        total = 0
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                if total > max_cells:
                    break
                for cell in row:
                    if cell is not None:
                        texts.append(str(cell))
                        total += 1
        return "\n".join(texts)
    except Exception:
        return ""

def extract_csv(path, max_rows=100000):
    try:
        texts = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i > max_rows:
                    break
                texts.append(" ".join([str(x) for x in row if x is not None]))
        return "\n".join(texts)
    except Exception:
        return ""

def extract_text(path, ext):
    ext = (ext or "").lower()
    if ext == ".pdf":
        return extract_pdf(path)
    elif ext == ".docx":
        return extract_docx(path)
    elif ext == ".pptx":
        return extract_pptx(path)
    elif ext == ".xlsx":
        return extract_xlsx(path)
    elif ext == ".txt":
        return safe_read_text(path)
    elif ext == ".csv":
        return extract_csv(path)
    else:
        return ""

def normalize_ext_list(ext: Optional[str]) -> Optional[List[str]]:
    if not ext:
        return None
    items = []
    for e in ext.split(","):
        e = e.strip().lower()
        if not e:
            continue
        if not e.startswith("."):
            e = "." + e
        items.append(e)
    return items or None

# --------- resposta ---------
class IndexRunResult(BaseModel):
    candidates: int
    indexed: int
    skipped: int
    errors: int
    elapsed_sec: float

# --------- endpoint ---------
@router.post("/run", response_model=IndexRunResult)
def index_run(
    db: Session = Depends(get_db),
    root_id: Optional[int] = Query(None, description="Se informado, indexa apenas essa raiz"),
    ext: Optional[str] = Query(None, description="Filtro por extensões: ex 'pdf,docx,xlsx'"),
    limit: Optional[int] = Query(None, ge=1, le=100000, description="Limite opcional de arquivos a processar"),
    reindex_all: bool = Query(False, description="Se true, força reindexação mesmo sem mudança")
):
    import time
    t0 = time.time()

    ext_filter = normalize_ext_list(ext)

    # Seleciona arquivos candidatos
    q = db.query(File)
    if root_id is not None:
        q = q.filter(File.root_id == root_id)
    # aplica filtro por extensão suportada/solicitada
    if ext_filter:
        q = q.filter(File.ext.in_(ext_filter))
    else:
        q = q.filter(File.ext.in_(list(SUPPORTED_EXTS)))

    # ordena por mtime desc e limita se pedido
    q = q.order_by(File.mtime.desc().nullslast(), File.id.asc())
    if limit:
        q = q.limit(limit)

    files = q.all()

    candidates = len(files)
    indexed = 0
    skipped = 0
    errors = 0

    # Preparar statements SQL para FTS5 (docs) e map
    ins_docs = text("INSERT INTO docs(content, filename, ext) VALUES (:content, :filename, :ext)")
    sel_map = text("SELECT rowid_docs, fingerprint FROM map WHERE file_id = :fid")
    upsert_map = text("""
        INSERT INTO map(file_id, rowid_docs, fingerprint) VALUES (:fid, :rowid, :fp)
        ON CONFLICT(file_id) DO UPDATE SET rowid_docs = excluded.rowid_docs, fingerprint = excluded.fingerprint
    """)
    del_docs = text("DELETE FROM docs WHERE rowid = :rowid")

    batch = 0

    for f in files:
        try:
            fp = f"{f.size}-{f.mtime}"
            # checar map existente
            row_map = db.execute(sel_map, {"fid": f.id}).fetchone()
            if row_map and not reindex_all:
                existing_rowid, existing_fp = row_map
                if existing_fp == fp:
                    skipped += 1
                    continue
                # fingerprint mudou -> apagar docs antigo antes de reindexar
                if existing_rowid is not None:
                    db.execute(del_docs, {"rowid": existing_rowid})

            # extrair texto
            content = extract_text(f.path, f.ext) or ""
            # inserir no FTS5
            res = db.execute(ins_docs, {"content": content, "filename": f.name, "ext": f.ext or ""})
            # SQLite retorna via cursor.lastrowid em nível DB-API; em SQLAlchemy, podemos buscar com SELECT last_insert_rowid()
            rowid = db.execute(text("SELECT last_insert_rowid()")).scalar()

            # mapear file_id -> docs.rowid com fingerprint
            db.execute(upsert_map, {"fid": f.id, "rowid": rowid, "fp": fp})

            indexed += 1
            batch += 1
            if batch % BATCH_SIZE == 0:
                db.commit()
        except Exception:
            errors += 1
            continue

    db.commit()
    dt = time.time() - t0

    return IndexRunResult(
        candidates=candidates,
        indexed=indexed,
        skipped=skipped,
        errors=errors,
        elapsed_sec=round(dt, 2),
    )
